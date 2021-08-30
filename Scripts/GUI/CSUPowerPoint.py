from pptx import Presentation
from pptx.util import Pt
from tkinter import *

#Defined as the functions needed to make a proper powerpoint for our meetings
class CSUPowerPoint:
    def __init__(self, tkinterInterface):
        #Set the Tkinter Interface Variable
        self.interface = tkinterInterface
        self.interface.title("CSU PowerPoint Generator")
        #The presentation uses the format/design of the chosen powerpoint
        self.presentation = Presentation("../../Files/format.pptx")
        #Has all officer information
        self.groupInfoFile = open("../../Files/groupInfo.txt", "r")
        
        #Start with introduction page
        self.titleSlide()
        self.officersSlide()
        self.introduction()

    #Function to simplify the calls to create a slide, add it to the presentation, and assign text to it
    def setSlideInfo(self, layout, titleText):
        #The layout of the slide (title slide, information slide, etc)
        slideLayout = self.presentation.slide_layouts[layout]
        #Add the slide in the specified layout
        slide = self.presentation.slides.add_slide(slideLayout)
        #Create a textbox and add the text specified
        textBox = slide.shapes.title
        textBox.text = titleText
        
        #Returns the slide layout, created slide and the created text box
        return (slideLayout, slide, textBox)
    
    #Returns the passage as a whole
    def getPassages(self, book, chapter, startVerse, endVerse):
        passage = []
        for lines in book:
            if not ":" in lines:
                continue
            #Gets the index of the colon
            firstColon = lines.index(":")
            #Assigns the current chapter number
            currentChapter = int(lines[0:firstColon])
            #Gets the second colon
            secondColon = lines.index(":", firstColon + 1,)
            #Assigns the current verse
            currentVerse = int(lines[firstColon+1:secondColon])
            #If the chapter equals the current chapter of the file and the verses match, then append the information
            if currentChapter == chapter and currentVerse >= startVerse and currentVerse <= endVerse:
                passage.append(lines[firstColon+1:])

        return passage

    #Check the passage for irregularities
    def checkPassage(self, userInput):
        #If the input starts with a digit for the book, then make sure it is set to "1 John" for example
        if userInput[0].isdigit():
            userInput = userInput[0] + " " + userInput[1].upper() + userInput[2:]
        elif userInput == "songofsolomon":#only book with spaces in it within the Bible
            userInput = "Song of Solomon"
        else:#Otherwise, set the first letter to a capital
            userInput = userInput[0].upper() + userInput[1:]

        return userInput

    #Performs the introduction page in the GUI
    def introduction(self):
        #Create the Introductory Frame
        self.introductoryFrame = Frame(self.interface)
        self.introductoryFrame.pack()
        #Creates an Introductory Message
        introductionLabel = Label(self.introductoryFrame, text="Welcome to the CSU PowerPoint Generator V1!\nClick the \"Next\" Button to continue!")
        introductionLabel.pack()
        #Adds a next button to move to the next page
        nextButton = Button(self.introductoryFrame, text="Next", fg="green", command=self.prayersAnnouncementsSlide)
        nextButton.pack()
        #Adds an exit button to close the application
        exitButton = Button(self.introductoryFrame, text="Exit", fg="red", command=self.interface.destroy)
        exitButton.pack()

    #Creates the title slide for the meetings
    def titleSlide(self):
        #Sets the slide layout, slide and textbox
        (slideLayout, slide, textBox) = self.setSlideInfo(0, self.groupInfoFile.readline())

        #Assigns the placeholder for the text box to assign more information to it
        textBox = slide.placeholders[1]
        textBox.text = self.groupInfoFile.readline() + "\n" + self.groupInfoFile.readline()

    #Displays all of the information of the officers of the organization in the groupInfoFile.txt
    def officersSlide(self):
        #Sets the slide layout, slide and textbox
        (slideLayout, slide, textBox) = self.setSlideInfo(1, self.groupInfoFile.readline())

        #Assigns a shape to the textbox
        textBox = slide.shapes[1]
        #Assigns a frame to the textbox
        textFrame = textBox.text_frame
        #Adds a paragraph
        lines = textFrame.paragraphs[0]
        #Reads the first line of the file
        lines.text = self.groupInfoFile.readline()
        lines.font.size = Pt(14)
        
        #Iterates through the file until it reaches the end
        person = self.groupInfoFile.readline()
        i=1
        while person != "":
            #Adds a paragraph for each person
            lines = textFrame.add_paragraph()
            #This is their information (Title: Name)
            lines.text = person
            lines.font.size = Pt(14)
            person = self.groupInfoFile.readline()
            i += 1

    def addPrayers(self, prayerFocus, textFrame):
        prayerFocus = prayerFocus.splitlines()
        for index in range(len(prayerFocus)):
            if(prayerFocus[index].startswith("  ")):
                lines = textFrame.add_paragraph()
                lines.text = prayerFocus[index]
                lines.level = 2
            else:
                lines = textFrame.add_paragraph()
                lines.text = prayerFocus[index]
                lines.level = 1
        
        self.announcementsFrame.destroy()
        self.otherAnnouncementsFrame.pack()

    def addOtherAnnouncements(self, announcements, textFrame):
        announcements = announcements.splitlines()
        for index in range(len(announcements)):
            if(announcements[index].startswith("  ")):
                lines = textFrame.add_paragraph()
                lines.text = announcements[index]
                lines.level = 1
            else:
                lines = textFrame.add_paragraph()
                lines.text = announcements[index]
                lines.level = 0
        
        self.otherAnnouncementsFrame.destroy()
        self.seriesAndLocationFrame.pack()

    def addSeriesAndLocation(self, series, location, textFrame):
        lines =  textFrame.add_paragraph()
        lines.text = "Series: " + series
        lines.level = 0

        lines = textFrame.add_paragraph()
        lines.text = "Meeting Place: " + location

        self.seriesAndLocationFrame.destroy()
        self.lessonSlides()

    def savePresentation(self):
        self.presentation.save("your_powerpoint.pptx")
        self.interface.destroy()

    #Creates the announcements slide
    def prayersAnnouncementsSlide(self):
        '''
            Section: Prayer Focus
        '''
        #Destroys all elements created within the introductory frame
        #The title and officers slides are generated based on the file input
        self.introductoryFrame.destroy()
        #Creates a new frame for the announcements input
        self.announcementsFrame = Frame(self.interface)
        self.announcementsFrame.pack()

        #Sets the slide layout, slide and textbox
        (slideLayout, slide, textBox) = self.setSlideInfo(1, "Announcements")
        #Adds a shape to the textbox
        textBox = slide.shapes[1]
        #Gets the text frame of the textbox
        textFrame = textBox.text_frame

        #Adds text to the frame
        lines = textFrame
        lines.text = "Prayer Focus"
        lines.level = 0
        
        #Prayer Focus (Section to be edited to work with GUI)
        announcementsLabel = Label(self.announcementsFrame, text="Enter prayer focuses below: ")
        announcementsLabel.pack()

        prayerFocusText = Text(self.announcementsFrame)
        prayerFocusText.pack()

        #Adds a next button to move to the next page
        nextButton = Button(self.announcementsFrame, text="Next", fg="green", command=lambda: self.addPrayers(prayerFocusText.get('1.0', 'end-1c'), textFrame))
        nextButton.pack()
        #Adds an exit button to close the application
        exitButton = Button(self.announcementsFrame, text="Exit", fg="red", command=self.savePresentation)
        exitButton.pack()

        '''
            Section: Other Announcements
        '''
        self.otherAnnouncementsFrame = Frame(self.interface)

        otherAnnouncementsLabel = Label(self.otherAnnouncementsFrame, text="Enter other announcements below:")
        otherAnnouncementsLabel.pack()

        otherAnnouncementsText = Text(self.otherAnnouncementsFrame)
        otherAnnouncementsText.pack()

        #Adds a next button to move to the next page
        nextButton = Button(self.otherAnnouncementsFrame, text="Next", fg="green", command=lambda: self.addOtherAnnouncements(otherAnnouncementsText.get('1.0', 'end-1c'), textFrame))
        nextButton.pack()
        #Adds an exit button to close the application
        exitButton = Button(self.otherAnnouncementsFrame, text="Exit", fg="red", command=self.savePresentation)
        exitButton.pack()

        '''
            Section: Series and Location
        '''
        self.seriesAndLocationFrame = Frame(self.interface)

        seriesLabel = Label(self.seriesAndLocationFrame, text="Enter Series:")
        seriesLabel.pack()

        seriesText = Text(self.seriesAndLocationFrame, height="10")
        seriesText.pack()

        locationLabel = Label(self.seriesAndLocationFrame, text="Enter Location:")
        locationLabel.pack()

        locationText = Text(self.seriesAndLocationFrame, height="10")
        locationText.pack()

        #Adds a next button to move to the next page
        nextButton = Button(self.seriesAndLocationFrame, text="Next", fg="green", command=lambda: self.addSeriesAndLocation(seriesText.get('1.0', 'end-1c'), locationText.get('1.0', 'end-1c'), textFrame))
        nextButton.pack()
        #Adds an exit button to close the application
        exitButton = Button(self.seriesAndLocationFrame, text="Exit", fg="red", command=self.savePresentation)
        exitButton.pack()
    
    '''
        Still need to do this portion!!!!
    '''
    #Creates a lesson header slide
    def lessonHeaderSlide(self):
        #Gets the title of the lesson from the user
        print("\nEnter the title of the lesson: ")
        userInput = self.getUserInput()

        #Sets the slide layout, slide and textbox
        slideLayout, slide, textBox = self.setSlideInfo(2, userInput)

        #Gets the name of the person teaching the lesson
        print("\nEnter the person teaching: ")
        userInput = self.getUserInput()
        textBox = slide.placeholders[1]
        textBox.text = userInput

    '''
        Figure out how to properly parse the selection
    '''
    #Creates a bible passage using the KJV folder (Helper Method)
    def biblePassage(self):
        #Destroy the lessonSlidesFrame
        self.lessonSlidesFrame.destroy()

        #Create Frame for biblePassage
        self.biblePassageFrame = Frame(self.interface, width="100")
        self.biblePassageFrame.pack()

        #Create a drop-down for selection of books from old and new testaments
        booksOfBible = [
            "Genesis",
            "Exodus",
            "Leviticus",
            "Numbers",
            "Deuteronomy",
            "Joshua",
            "Judges",
            "Ruth",
            "1 Samuel",
            "2 Samuel",
            "1 Kings",
            "2 Kings",
            "1 Chronicles",
            "2 Chronicles",
            "Ezra",
            "Job",
            "Psalms",
            "Proverbs",
            "Ecclesiastes",
            "Song of Solomon",
            "Isaiah",
            "Jeremiah",
            "Lamentations",
            "Ezekiel",
            "Daniel",
            "Hosea",
            "Joel",
            "Amos",
            "Obadiah",
            "Jonah",
            "Micah",
            "Nahum",
            "Habakkuk",
            "Zephaniah",
            "Haggai",
            "Zechariah",
            "Malachi",
            "Matthew",
            "Mark",
            "Luke",
            "John",
            "Acts",
            "Romans",
            "1 Corinthians",
            "2 Corinthians",
            "Galatians",
            "Ephesians",
            "Philippians",
            "Colossians",
            "1 Thessalonians",
            "2 Thessalonians",
            "1 Timothy",
            "2 Timothy",
            "Titus",
            "Philemon",
            "Hebrews",
            "James",
            "1 Peter",
            "2 Peter",
            "1 John",
            "2 John",
            "3 John",
            "Jude",
            "Revelation"
        ]
        selectedBook = StringVar(self.biblePassageFrame)
        selectedBook.set("Select a Book")

        #This does not work well....need to see if we can change how it works or if there are other options
        #Online says there is another library called ttk or something like that
        bibleBookDropDown = OptionMenu(self.biblePassageFrame, selectedBook, *booksOfBible)
        bibleBookDropDown.config(width=15)
        bibleBookDropDown.pack()

        #Get the chapter from the user
        chapterLabel = Label(self.biblePassageFrame, text="Enter the Chapter: ")
        chapterLabel.pack()
        chapterText = Text(self.biblePassageFrame, height=1)
        chapterText.pack()

        #Get the start verse
        verseStartLabel = Label(self.biblePassageFrame, text="Enter the Start Verse: ")
        verseStartLabel.pack()
        verseStart = Text(self.biblePassageFrame, height=1)
        verseStart.pack()

        #Get the end verse
        verseEndLabel = Label(self.biblePassageFrame, text="Enter the end verse: ")
        verseEndLabel.pack()
        verseEnd = Text(self.biblePassageFrame, height=1)
        verseEnd.pack()

        #Adds a next button to move to the next page
        #selectedBook = str(selectedBook.get().replace(" ", "").lower())
        continueButton = Button(self.biblePassageFrame, text="Continue", fg="green", command=lambda: self)
        continueButton.pack()

        '''
        try:
            #Gets the book from th euser
            print("\nEnter book of the Bible: ")
            #Sets it to lowercase after replacing the spaces with empty strings
            userInput = input("> ").replace(" ", "").lower()
            #Grabs the book text file
            book = open('../Files/KJV/' + userInput + '.txt', 'r')
            #Asks the user for the chapter
            chapter = self.getChapterAndVerses("Chapter: ")
            #Asks the user for the starting verse
            startVerse = self.getChapterAndVerses("Starting Verse: ")
            #Asks the user for the ending verse
            endVerse = self.getChapterAndVerses("Ending Verse: ")

            #If the end verse is less than the start verse, raise an error
            if endVerse < startVerse:
                raise(LookupError)
            
            #Iterate through the specified book to find the chapter
            passage = self.getPassages(book, chapter, startVerse, endVerse)

            #If the passage length is 0, raise a lookup error
            if(len(passage) == 0):
                raise(LookupError)
            
            #Check for digits (1 John), or spaces (Song of Solomon)
            userInput = self.checkPassage(userInput)

            #If they are equal, then we only put one verse
            if(startVerse == endVerse):
                title = userInput + " " + str(chapter) + ":" + str(startVerse)
            else:#if they are not, then we have multiple verses denoted as (Book Chapter:StartVerse - Endverse)
                title = userInput + " " + str(chapter) + ":" + str(startVerse) + "-" + passage[len(passage) - 1][:passage[len(passage)-1].index(":")]

            #Iterate through the passage with a maximum of 9 lines per slide
            linesUsed = 9
            for lines in passage:
                #Gets the character in the verse
                charInVerse = len(lines) / 50 + 1
                #If it is greater than 9, then create a new slide to add text
                if(charInVerse + linesUsed > 9):
                    (slideLayout, slide, textBox) = self.setSlideInfo(1, title)
                    textBox = slide.shapes.placeholders[1]
                    textBox.text = lines
                    textFrame = textBox.text_frame
                    linesUsed = 1
                    continue
                #Add the verse to the frame
                verse = textFrame.add_paragraph()
                verse.text = lines
                linesUsed += charInVerse
        except OSError:#If the user enters an invalid book, tell them it is invalid
            print("Invalid Book")
        except ValueError:#If the user enters an invalid number, tell them it is invalid
            print("Must enter a valid number")
        except LookupError:#If the user enters an invalid verse, tell them it is invalid
            print("Invalid verses. No such passage exists")
        '''

    #Creates the header slide of the presentation (Helper Method)
    def header(self):
        #Destroy the lessonSlidesFrame
        self.lessonSlidesFrame.destroy()

        #Create the header frame
        self.headerSlideFrame = Frame(self.interface)
        self.headerSlideFrame.pack()

        headerLabel = Label(self.headerSlideFrame, text="Enter the header: ")
        headerLabel.pack()
        headerText = Text(self.headerSlideFrame, height=1)
        headerText.pack()

        continueButton = Button(self.headerSlideFrame, text="Continue", fg="green", command=lambda: self.headerToLessonSlides(headerText.get('1.0', 'end-1c')))
        continueButton.pack()

    def headerToLessonSlides(self, title):
        #Assigns the values
        (slideLayout, slide, textBox) = self.setSlideInfo(2, title)
        self.headerSlideFrame.destroy()
        self.lessonSlides()

    def addPoints(self, title, points):
        #Creates the slide and textbox
        (slideLayout, slide, textBox) = self.setSlideInfo(1, title)
        #Assigns a placeholder for the frame
        textBox = slide.shapes.placeholders[1]
        textFrame = textBox.text_frame

        points = points.splitlines()
        for index in range(len(points)):
            if(points[index].startswith("  ")):
                lines = textFrame.add_paragraph()
                lines.text = points[index]
                lines.level = 1
            else:
                lines = textFrame.add_paragraph()
                lines.text = points[index]
                lines.level = 0
        
        self.pointsSlideFrame.destroy()
        self.lessonSlides()

    #Defines the points of the message (Helper Method)
    def points(self):
        #Destroy the lessonSlidesFrame
        self.lessonSlidesFrame.destroy()

        #Create PointsSlideFrame
        self.pointsSlideFrame = Frame(self.interface)
        self.pointsSlideFrame.pack()

        pointsTitleLabel = Label(self.pointsSlideFrame, text="Enter Title of Slide: ")
        pointsTitleLabel.pack()
        pointsTitleText = Text(self.pointsSlideFrame, height=1)
        pointsTitleText.pack()

        pointsLabel = Label(self.pointsSlideFrame, text="Enter points below: ")
        pointsLabel.pack()
        pointsText = Text(self.pointsSlideFrame)
        pointsText.pack()

        continueButton = Button(self.pointsSlideFrame, text="Continue", fg="green", command=lambda: self.addPoints(pointsTitleText.get('1.0', 'end-1c'), pointsText.get('1.0', 'end-1c')))
        continueButton.pack()

    #Defines the lesson slides and asks the user for input each time until they quit
    def lessonSlides(self):
        #Add Lesson Slides frame
        self.lessonSlidesFrame = Frame(self.interface)
        self.lessonSlidesFrame.pack()

        biblePassageButton = Button(self.lessonSlidesFrame, text="Add Bible Passage", command=lambda: self.biblePassage())
        biblePassageButton.pack()
        headerButton = Button(self.lessonSlidesFrame, text="Add Header", command=lambda: self.header())
        headerButton.pack()
        pointsButton = Button(self.lessonSlidesFrame, text="Add Points", command=lambda: self.points())
        pointsButton.pack()

        #Adds a next button to move to the next page
        saveButton = Button(self.lessonSlidesFrame, text="Save Presentation", fg="green", command=lambda: self.savePresentation())
        saveButton.pack()