from pptx import Presentation
from pptx.util import Pt

class CSUPowerPoint:
    def __init__(self):
        self.presentation = Presentation("../Files/format.pptx")
        self.groupInfoFile = open("../Files/groupInfo.txt", "r")

    def introduction(self):
        print("\nWelcome to the CSU PowerPoint Generator V1!")

    #Does this work????, If so, can it be rewritten where it is used?
    #Do we need slideLayout for anything other than the textBox??????
    #What happens if we don't have it in each method that uses it?
    def setSlideInfo(self, layout, titleText):
        slideLayout = self.presentation.slide_layouts[layout]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = titleText
        
        return (slideLayout, slide, textBox)

    def titleSlide(self):
        #Repeated
        """
        slideLayout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = self.groupInfoFile.readline()
        """
        #Repeated
        slideLayout, slide, textBox = self.setSlideInfo(0, self.groupInfoFile.readline())

        textBox = slide.placeholders[1]
        textBox.text = self.groupInfoFile.readline() + "\n" + self.groupInfoFile.readline()

    def officersSlide(self):
        #Repeated
        """
        slideLayout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = self.groupInfoFile.readline()
        """
        #Repeated
        slideLayout, slide, textBox = self.setSlideInfo(1, self.groupInfoFile.readline())

        textBox = slide.shapes[1]
        textFrame = textBox.text_frame
        lines = textFrame.paragraphs[0]
        lines.text = self.groupInfoFile.readline()
        lines.font.size = Pt(14)
        
        person = self.groupInfoFile.readline()
        i=1
        while person != "":
            lines = textFrame.add_paragraph()
            lines.text = person
            lines.font.size = Pt(14)
            person = self.groupInfoFile.readline()
            i += 1

    def announcementsSlide(self):
        #Repeated
        """
        slideLayout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = "Announcements"
        """
        #Repeated
        slideLayout, slide, textBox = self.setSlideInfo(1, "Announcements")

        textBox = slide.shapes[1]
        textFrame = textBox.text_frame
        lines = textFrame
        lines.text = "Prayer Focus"
        lines.level = 0
        
        #Prayer Focus
        print("\nPrayer Focus List: ")
        print("Enter each prayer focus and then hit enter.")
        print("After the last focus, type \'done\' and hit enter to move on: ")
        while True:
            lines = textFrame.add_paragraph()
            userInput = input("> ")
            while len(userInput) == 0:
                userInput = input("> ")
            if userInput.lower().replace(" ", "") == "done":
                break
            lines.text = userInput
            lines.level = 1

        #Current Series
        lines = textFrame.add_paragraph()
        print("Enter the name of the series you are going through: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        lines.text = "Series: " + userInput
        
        #Location of Meetings
        lines = textFrame.add_paragraph()
        lines = textFrame.add_paragraph()
        print("\nEnter where you are holding meetings: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        lines.text = "Meeting place: " + userInput
    
    def lessonHeaderSlide(self):
        print("\nEnter the title of the lesson: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        #Repeated
        """
        slideLayout = self.presentation.slide_layouts[2]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = userInput
        """
        #Repeated
        slideLayout, slide, textBox = self.setSlideInfo(2, userInput)

        print("\nEnter the person teaching: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        textBox = slide.placeholders[1]
        textBox.text = userInput
    
    def biblePassage(self):
        try:
            print("\nEnter book of the Bible: ")
            userInput = input("> ").replace(" ", "").lower()
            book = open('../Files/KJV/' + userInput + '.txt', 'r')
            print("Chapter: ")
            chapter = int(input("> "))
            print("Starting Verse: ")
            startVerse = int(input("> "))
            print("Ending Verse: ")
            endVerse = int(input("> "))

            if endVerse < startVerse:
                raise(LookupError)
            
            passage = []
            for lines in book:
                if not ":" in lines:
                    continue
                firstColon = lines.index(":")
                currentChapter = int(lines[0:firstColon])
                secondColon = lines.index(":", firstColon + 1,)
                currentVerse = int(lines[firstColon+1:secondColon])
                if currentChapter == chapter and currentVerse >= startVerse and currentVerse <= endVerse:
                    passage.append(lines[firstColon+1:])

            if(len(passage) == 0):
                raise(LookupError)
            
            if userInput[0].isdigit():
                userInput = userInput[0] + " " + userInput[1].upper() + userInput[2:]
            elif userInput == "songofsolomon":
                userInput = "Song of Solomon"
            else:
                userInput = userInput[0].upper() + userInput[1:]

            if(startVerse == endVerse):
                title = userInput + " " + str(chapter) + ":" + str(startVerse)
            else:
                title = userInput + " " + str(chapter) + ":" + str(startVerse) + "-" + passage[len(passage) - 1][:passage[len(passage)-1].index(":")]

            linesUsed = 9
            for lines in passage:
                charInVerse = len(lines) / 50 + 1
                if(charInVerse + linesUsed > 9):
                    #Repeated
                    slideLayout = self.presentation.slide_layouts[1]
                    slide = self.presentation.slides.add_slide(slideLayout)
                    textBox = slide.shapes.title
                    textBox.text = title
                    #Repeated
                    textBox = slide.shapes.placeholders[1]
                    textBox.text = lines
                    textFrame = textBox.text_frame
                    linesUsed = 1
                    continue
                verse = textFrame.add_paragraph()
                verse.text = lines
                linesUsed += charInVerse
        except OSError:
            print("Invalid Book")
        except ValueError:
            print("Must enter a valid number")
        except LookupError:
            print("Invalid verses. No such passage exists")

    def header(self):
        print("Enter the header: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        #Repeated
        slideLayout = self.presentation.slide_layouts[2]
        slide = self.presentation.slides.add_slide(slideLayout)
        textBox = slide.shapes.title
        textBox.text = userInput
        #Repeated

    def points(self):
        slideLayout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slideLayout)
        print("Title of slide: ")
        userInput = input("> ")
        while len(userInput) == 0:
            userInput = input("> ")
        textBox = slide.shapes.title
        textBox.text = userInput
        textBox = slide.shapes.placeholders[1]
        textFrame = textBox.text_frame
        print("Enter the points for the slide. Type \'done\' to finish slide")
        textFrame.text = input("> ") + '\n'
        while True:
            userInput = input("> ")
            while userInput == 0:
                userInput = input("> ")
            if userInput.lower().replace(" ", "") == 'done':
                break
            lines = textFrame.add_paragraph()
            lines.text = userInput + '\n'

    def lessonSlides(self):
        while True:
            print("\nChoices: \'B\' for Bible Passage, \'H\' for slide with a header, \'P\' for slide with points, \'quit\' to complete lesson")
            userInput = input("> ")
            userInput = userInput.upper().replace(" ", "")
            if(userInput == 'QUIT'):
                break
            elif(userInput == 'B'):
                self.biblePassage()
            elif(userInput == 'H'):
                self.header()
            elif(userInput == 'P'):
                self.points()
            else:
                print("Enter a valid choice")

        self.presentation.save("your_powerpoint.pptx")
