# this program is meant to automatically create decent looking powerpoints for CSU meetings
# Written by Nicolas Null
# Spring 2021

# all information like announcements, leaders and such will be kept in the text file groupInfo.txt
# uses library python-pptx
from pptx import Presentation
from pptx.util import Pt

print("\nWelcome to the CSU PowerPoint Generator V1!")
prs = Presentation("format.pptx")
group_info_file = open("groupInfo.txt", "r")


# title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
text_box=slide.shapes.title
text_box.text=group_info_file.readline()
text_box=slide.placeholders[1]
text_box.text=group_info_file.readline()+"\n"+group_info_file.readline()

# Officers Slide
slide_layout = prs.slide_layouts[1]
slide=prs.slides.add_slide(slide_layout)
text_box=slide.shapes.title
text_box.text=group_info_file.readline()

text_box = slide.shapes[1]
text_frame = text_box.text_frame
lines = text_frame.paragraphs[0]
lines.text=group_info_file.readline()
lines.font.size= Pt(14)
s = group_info_file.readline()
i=1
while s != "":
    lines=text_frame.add_paragraph()
    lines.text=s
    lines.font.size= Pt(14)
    s=group_info_file.readline()
    i+=1

# announcements Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
text_box=slide.shapes.title
text_box.text="Announcements"

text_box=slide.shapes[1]
text_frame = text_box.text_frame
lines = text_frame
lines.text = "Prayer Focus"
lines.level = 0

print("\nPrayer Focus List:")
print("Enter each prayer focus and then hit enter.")
print("After the last focus type \'done\' and hit enter to move on:")
while True:
    lines=text_frame.add_paragraph()
    s = input("> ")
    while len(s) == 0:
        s = input("> ")
    if s.lower().replace(" ", "") == "done": 
        break
    lines.text= s
    lines.level=1


lines =text_frame.add_paragraph()
print("\nEnter the name of the series you are going through: ")
s = input("> ")
while len(s) == 0:
    s = input("> ")
lines.text="Series: " + s

lines=text_frame.add_paragraph()

lines=text_frame.add_paragraph()
print("\nEnter where you are holding meetings: ")
s = input("> ")
while len(s) == 0:
    s = input("> ")
lines.text = "Meetingplace: " + s


# lesson header slide
slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(slide_layout)
print("\nLesson:")
print("Enter the title of the lesson: ")
s = input("> ")
while len(s) == 0:
    s = input("> ")
text_box = slide.shapes.title
text_box.text = s

print("\nEnter the person teaching: ")
s = input("> ")
while len(s) == 0:
    s = input("> ")
text_box = slide.placeholders[1]
text_box.text=s

# now need to put lesson together with a combination of headers, points and bible passages

# bible_passage, header and points methods used to create that type of slide/slides in lesson section of PowerPoint
def bible_passage():
    try:
        print("\nBook of the Bible: ")
        s = input("> ").replace(" ", "").lower()
        book = open('KJV/'+s+".txt",'r')# makes sure to remove whitespace so file names will match up
        print("Chapter: ")
        chapter = int(input("> "))
        print("Starting Verse: ")
        start_verse = int(input("> "))
        print("Ending Verse: ")
        end_verse = int(input("> "))

        if end_verse < start_verse:# makes sure user doesn't enter an inverted range of verses
            raise(LookupError)

        # creates a string list, each index being a verse within the requested passage
        passage = []
        for lines in book:
            if not ":" in lines:#which means line does not contain a verse (name of book probably)
                continue
            first_colon = lines.index(":")
            current_chapter = int(lines[0:first_colon])
            second_colon = lines.index(":",first_colon+1,)
            current_verse = int(lines[first_colon+1:second_colon])
            if current_chapter == chapter and current_verse >= start_verse and current_verse <= end_verse:
                passage.append(lines[first_colon+1:])

        # this makes sure the verse range given is valid
        if(len(passage) == 0):
            raise(LookupError)

        # formats the title from the user input
        # need to capitalize all of the first letters
        # all books of bible are either 1 word or a number and a word. 
        # the only exception is Song of Solomon
        if(s[0].isdigit()): # add single space between number and first letter and capitalize first letter
            s = s[0] + " " + s[1].upper() + s[2:]
        elif s == "songofsolomon": # song of solomon only book with multiple words
            s = "Song of Solomon"
        else: # simply a one word book
            s = s[0].upper() + s[1:]

        if(start_verse == end_verse):
            title = s+" "+str(chapter)+":"+str(start_verse)
        else: # this checks that last verse is actually last verse not the verse user inputted
            title = s+" "+str(chapter)+":"+str(start_verse)+"-"+passage[len(passage)-1][:passage[len(passage)-1].index(":")]

        lines_used = 9
        for lines in passage:
            # safe number is 50 chars per line and 9 lines per slide
            char_in_verse = len(lines)/50 + 1
            if(char_in_verse+lines_used > 9):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                text_box = slide.shapes.title
                text_box.text = title
                text_box = slide.shapes.placeholders[1]
                text_box.text = lines
                text_frame = text_box.text_frame
                lines_used = 1
                continue
            verse = text_frame.add_paragraph()
            verse.text = lines
            lines_used += char_in_verse
            
                
    except OSError:
        print("Invalid book.")
    except ValueError:
        print("Must enter a valid number.")
    except LookupError:
        print("Invalid verses. No such passage exists.")



def header():
    print("Enter the header: ")
    s = input("> ")
    while len(s) == 0:
        s = input("> ")
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    text_box = slide.shapes.title
    text_box.text = s

def points():
    slide_layout=prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    print("Title of slide: ")
    s = input("> ")
    while len(s) == 0:
        s = input("> ")
    text_box = slide.shapes.title
    text_box.text = s
    text_box = slide.shapes.placeholders[1]
    text_frame = text_box.text_frame
    print("Enter the points for the slide. Type \'done\' to finish slide")
    text_frame.text = input("> ") + '\n'
    while True:
        s = input("> ")
        while len(s) == 0:
            s = input("> ")
        if s.lower().replace(" ","") == 'done':
            break
        lines = text_frame.add_paragraph()
        lines.text = s + '\n'

while True: # this loops through the above methods until done with lesson
    print("\nChoices: \'B\' for bible passage, \'H\' for slide with a header, \'P\' for slide with points, \'quit\' to complete lesson")
    s = input("> ")
    s=s.upper().replace(" ","")
    if(s =='QUIT'):
        break
    elif(s=='B'):
        bible_passage()
    elif(s=='H'):
        header()
    elif(s=='P'):
        points()
    else:
        print("enter a valid choice")

# saves powerpoint under the name your_powerpoint.pptx    
prs.save("your_powerpoint.pptx")